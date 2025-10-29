"""Document content extraction."""

from pathlib import Path
from typing import Optional

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None


class DocumentExtractor:
    """Extract text content from various document formats."""

    def __init__(self, max_length: int = 2000, max_pages: int = 5):
        """
        Initialize document extractor.

        Args:
            max_length: Maximum characters to extract
            max_pages: Maximum pages to process (for PDFs)
        """
        self.max_length = max_length
        self.max_pages = max_pages

    def extract(self, file_path: Path) -> Optional[str]:
        """
        Extract text from document.

        Args:
            file_path: Path to document file

        Returns:
            Extracted text or None on error
        """
        extension = file_path.suffix.lower()

        try:
            if extension == ".pdf":
                return self._extract_pdf(file_path)
            elif extension in [".docx", ".doc"]:
                return self._extract_docx(file_path)
            elif extension in [".pptx", ".ppt"]:
                return self._extract_pptx(file_path)
            elif extension in [".xlsx", ".xls"]:
                return self._extract_xlsx(file_path)
            elif extension in [".txt", ".srt", ".md", ".csv"]:
                return self._extract_txt(file_path)
            else:
                return None
        except Exception as e:
            print(f"Error extracting from {file_path}: {e}")
            return None

    def _extract_pdf(self, file_path: Path) -> Optional[str]:
        """Extract text from PDF."""
        if PyPDF2 is None:
            print("PyPDF2 not installed. Install with: pip install PyPDF2")
            return None

        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text_parts = []

                # Limit number of pages
                num_pages = min(len(reader.pages), self.max_pages)

                for i in range(num_pages):
                    page = reader.pages[i]
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)

                    # Check if we've extracted enough
                    current_length = sum(len(part) for part in text_parts)
                    if current_length >= self.max_length:
                        break

                full_text = "\n".join(text_parts)
                return full_text[: self.max_length]

        except Exception as e:
            print(f"Error reading PDF {file_path}: {e}")
            return None

    def _extract_docx(self, file_path: Path) -> Optional[str]:
        """Extract text from DOCX."""
        if docx is None:
            print("python-docx not installed. Install with: pip install python-docx")
            return None

        try:
            doc = docx.Document(file_path)
            text_parts = []

            for paragraph in doc.paragraphs:
                text_parts.append(paragraph.text)

                # Check length
                current_length = sum(len(part) for part in text_parts)
                if current_length >= self.max_length:
                    break

            full_text = "\n".join(text_parts)
            return full_text[: self.max_length]

        except Exception as e:
            print(f"Error reading DOCX {file_path}: {e}")
            return None

    def _extract_pptx(self, file_path: Path) -> Optional[str]:
        """Extract text from PPTX."""
        if Presentation is None:
            print("python-pptx not installed. Install with: pip install python-pptx")
            return None

        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)

                # Check length
                current_length = sum(len(part) for part in text_parts)
                if current_length >= self.max_length:
                    break

            full_text = "\n".join(text_parts)
            return full_text[: self.max_length]

        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}")
            return None

    def _extract_txt(self, file_path: Path) -> Optional[str]:
        """Extract text from TXT file."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read(self.max_length)
                return text
        except Exception as e:
            print(f"Error reading TXT {file_path}: {e}")
            return None

    def _extract_xlsx(self, file_path: Path) -> Optional[str]:
        """Extract text from Excel file."""
        if openpyxl is None:
            print("openpyxl not installed. Install with: pip install openpyxl")
            return None

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheet = wb.active

            # Get sheet name
            summary = f"Sheet: {sheet.title}\n\n"

            # Get first row (headers)
            headers = []
            for cell in list(sheet.rows)[0]:
                if cell.value:
                    headers.append(str(cell.value))

            if headers:
                summary += f"Columns: {', '.join(headers[:10])}\n\n"

            # Get first few rows of data
            summary += "Sample data:\n"
            for i, row in enumerate(list(sheet.rows)[1:6], 1):  # First 5 data rows
                row_data = [str(cell.value) if cell.value else "" for cell in row[:5]]
                if any(row_data):
                    summary += f"Row {i}: {' | '.join(row_data)}\n"

            return summary[: self.max_length]

        except Exception as e:
            print(f"Error reading Excel {file_path}: {e}")
            return None
